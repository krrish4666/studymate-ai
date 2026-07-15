import io
from pathlib import Path

from app.core.exceptions import BadRequestError


class FileParser:
    def extract_text(self, file_data: bytes, filename: str) -> str:
        ext = Path(filename).suffix.lower()

        if ext == ".pdf":
            return self._parse_pdf(file_data)
        elif ext == ".docx":
            return self._parse_docx(file_data)
        elif ext == ".pptx":
            return self._parse_pptx(file_data)
        elif ext == ".txt":
            return file_data.decode("utf-8", errors="replace")
        elif ext in (".jpg", ".jpeg", ".png", ".webp"):
            return self._parse_image(file_data)
        else:
            raise BadRequestError(f"Unsupported file type: {ext}")

    def _parse_pdf(self, data: bytes) -> str:
        import fitz
        doc = fitz.open(stream=data, filetype="pdf")
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return text.strip()

    def _parse_docx(self, data: bytes) -> str:
        import docx
        doc = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs).strip()

    def _parse_pptx(self, data: bytes) -> str:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts).strip()

    def _parse_image(self, data: bytes) -> str:
        from PIL import Image
        import base64
        img = Image.open(io.BytesIO(data))
        if img.mode != "RGB":
            img = img.convert("RGB")
        max_size = (1024, 1024)
        img.thumbnail(max_size, Image.LANCZOS)
        buffer = io.BytesIO()
        img.save(buffer, format="JPEG", quality=85)
        return base64.b64encode(buffer.getvalue()).decode("utf-8")


file_parser = FileParser()
