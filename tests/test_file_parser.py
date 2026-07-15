import pytest

from app.services.file_parser import file_parser
from app.core.exceptions import BadRequestError


class TestFileParser:
    def test_parse_txt(self):
        text = file_parser.extract_text(b"Hello, World!", "test.txt")
        assert text == "Hello, World!"

    def test_parse_txt_unicode(self):
        text = file_parser.extract_text("Café résumé".encode("utf-8"), "notes.txt")
        assert "Café" in text
        assert "résumé" in text

    def test_parse_txt_multiline(self):
        content = b"Line 1\nLine 2\nLine 3"
        text = file_parser.extract_text(content, "doc.txt")
        assert "Line 1" in text
        assert "Line 2" in text
        assert "Line 3" in text

    def test_parse_pdf(self):
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "PDF Content Test", fontsize=12)
        pdf_bytes = doc.write()
        doc.close()

        text = file_parser.extract_text(pdf_bytes, "test.pdf")
        assert "PDF Content Test" in text

    def test_parse_docx(self):
        import docx
        doc = docx.Document()
        doc.add_paragraph("DOCX Content Test")
        doc.add_paragraph("Second paragraph")
        import io
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)

        text = file_parser.extract_text(buf.read(), "test.docx")
        assert "DOCX Content Test" in text
        assert "Second paragraph" in text

    def test_parse_pptx(self):
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "PPTX Title"
        import io
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        text = file_parser.extract_text(buf.read(), "test.pptx")
        assert "PPTX Title" in text

    def test_parse_image_returns_base64(self):
        from PIL import Image
        import io
        img = Image.new("RGB", (100, 100), color="red")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)

        result = file_parser.extract_text(buf.read(), "test.png")
        assert isinstance(result, str)
        assert len(result) > 100

    def test_unsupported_extension(self):
        with pytest.raises(BadRequestError, match="Unsupported file type"):
            file_parser.extract_text(b"some data", "file.xyz")
